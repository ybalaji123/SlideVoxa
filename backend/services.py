import os
import json
import base64
import aiohttp
import asyncio
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv

load_dotenv(os.path.join(os.path.dirname(__file__), "..", "..", ".env"))

import sys
import io

def safe_str(s):
    """Sanitize string for safe printing on Windows console."""
    if s is None: return ""
    try:
        # Get string representation
        val = str(s)
        # Encode to bytes using 'ascii' and 'replace' to handle non-ascii safely
        # Then decode back to string. This is the safest way to prevent charmap errors.
        return val.encode('ascii', 'replace').decode('ascii')
    except Exception:
        return "[sanitization failed]"

SARVAM_API_KEY = os.getenv("SARVAM_API_KEY", "").strip()
ELEVENLABS_API_KEY = os.getenv("Eleven_Labs", "").strip()

print(f"[services] SARVAM key loaded: {'YES' if SARVAM_API_KEY else 'MISSING'}")
print(f"[services] ElevenLabs key loaded: {'YES' if ELEVENLABS_API_KEY else 'MISSING'}")


import io
import tempfile
import shutil

def fast_extract_text(content: bytes) -> list[dict]:
    """Extremely fast text extraction using python-pptx (milliseconds)."""
    try:
        b = io.BytesIO(content)
        prs = Presentation(b)
        slides_data = []
        for i, slide in enumerate(prs.slides, 1):
            title = ""
            points = []
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                text = shape.text_frame.text.strip()
                if not text: continue
                if getattr(shape, "is_placeholder", False) and shape.placeholder_format.idx in (0, 1):
                    title = text
                    continue
                for para in shape.text_frame.paragraphs:
                    if para.text.strip(): points.append(para.text.strip())
            
            slides_data.append({
                "slide_number": i,
                "title": title or f"Slide {i}",
                "points": list(dict.fromkeys(points))[:10], # Limit points for speed
                "body": "", # Filled later if needed
                "image_data_uri": None
            })
        return slides_data
    except Exception as e:
        print(f"[fast_extract] Failed: {safe_str(e)}")
        return []

def extract_slides(content: bytes, presentation_id: str) -> list[dict]:
    """Extract slide components using stable slide-by-slide export for maximum reliability."""
    slide_image_b64 = {}
    temp_dir = tempfile.mkdtemp(prefix="slidevoxa_")
    temp_pptx = os.path.join(temp_dir, "temp_pres.pptx")
    
    with open(temp_pptx, "wb") as f:
        f.write(content)

    print(f"[extract] Starting extraction for {presentation_id} at {safe_str(temp_pptx)}")

    try:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        
        # Use a fresh instance of PowerPoint
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        powerpoint.Visible = False # Run invisibly
        
        abs_temp_pptx = os.path.abspath(temp_pptx)
        
        # Open presentation
        # Open(Filename, ReadOnly, Untitled, WithWindow)
        presentation = powerpoint.Presentations.Open(abs_temp_pptx, ReadOnly=1, WithWindow=False)
        
        total_slides = presentation.Slides.Count
        print(f"[extract] PowerPoint slide count: {total_slides}")
        
        # Get dimensions for crisp PNG export (1.5x to 2x scale for beauty)
        slide_width_pts = presentation.PageSetup.SlideWidth
        slide_height_pts = presentation.PageSetup.SlideHeight
        export_width = 1600 # Premium resolution
        export_height = int(export_width * slide_height_pts / slide_width_pts)
        
        for i, slide in enumerate(presentation.Slides, 1):
            img_path = os.path.join(temp_dir, f"s_{i}.png")
            # Export(FileName, FilterName, ScaleWidth, ScaleHeight)
            slide.Export(img_path, "PNG", export_width, export_height)
            
            if os.path.exists(img_path):
                with open(img_path, "rb") as im_f:
                    im_bytes = im_f.read()
                slide_image_b64[i] = "data:image/png;base64," + base64.b64encode(im_bytes).decode("utf-8")
                if i % 5 == 0: print(f"[extract] Exported {i}/{total_slides} slides...")
            else:
                print(f"[extract] WARNING: Failed to export slide {i}")

        presentation.Close()
        # powerpoint.Quit() # Don't quit if others are using it, but we used DispatchEx so we should
        try: powerpoint.Quit()
        except: pass
        
        print(f"[extract] Successfully exported {len(slide_image_b64)} images")
    except Exception as e:
        import traceback
        print(f"[extract] FAILED: {safe_str(e)}")
        print(f"[extract] {safe_str(traceback.format_exc())}")
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
        try: pythoncom.CoUninitialize()
        except: pass

    # Combine with text extraction
    prs = Presentation(io.BytesIO(content))
    slides_data = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        points = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    if getattr(shape, "is_placeholder", False) and shape.placeholder_format.idx in (0,1):
                        title = text
                    else:
                        points.append(text)
        
        slides_data.append({
            "slide_number": i,
            "title": title or f"Slide {i}",
            "points": list(dict.fromkeys(points)),
            "image_data_uri": slide_image_b64.get(i),
        })
    return slides_data


async def async_generate_script_sarvam(slide: dict, session: aiohttp.ClientSession) -> str:
    """Generate a natural speaking script for a slide using Sarvam AI asynchronously."""
    title = slide.get("title", "")
    points = slide.get("points", [])
    body = slide.get("body", "")

    bullet_text = "\n".join(f"- {p}" for p in points) if points else ""
    content = f"Title: {title}\n{bullet_text}\n{body}".strip()

    prompt = (
        "You are a professional AI presenter. Convert the following slide content "
        "into a natural, engaging 2-3 sentence spoken presentation script. "
        "Sound like a human presenter — clear, confident, and engaging. "
        "Do not use markdown. Just plain spoken text.\n\n"
        f"Slide content:\n{content}\n\nScript:"
    )

    try:
        async with session.post(
            "https://api.sarvam.ai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {SARVAM_API_KEY}",
                "Content-Type": "application/json",
            },
            json={
                "model": "sarvam-m",
                "messages": [{"role": "user", "content": prompt}],
                "max_tokens": 300,
            },
            timeout=30,
        ) as response:
            if response.status == 200:
                data = await response.json()
                content = data["choices"][0]["message"]["content"].strip()
                import re
                # Sarvam AI sometimes uses DeepSeek R1 and outputs <think> debug blocks; strip them out
                content = re.sub(r'<think>.*?</think>', '', content, flags=re.DOTALL).strip()
                return content
            else:
                text = await response.text()
                print(f"[Sarvam script] Error {response.status}: {safe_str(text[:200])}")
    except Exception as e:
        print(f"[Sarvam script] Exception: {safe_str(e)}")

    # Fallback: build a basic script from slide content
    script = f"In this slide, we discuss {title}."
    if points:
        script += " The key points covered are: " + "; ".join(points[:4]) + "."
    if body:
        script += " " + body[:200]
    return script


async def async_generate_audio_for_slide(script: str, slide_number: int, presentation_id: str, voice_id: str, session: aiohttp.ClientSession) -> str | None:
    """
    Generate audio for a slide asynchronously. Tries ElevenLabs first, falls back to Sarvam.
    Returns a base64-encoded audio string (data URI) or None.
    """
    # Try ElevenLabs first if key is available
    if ELEVENLABS_API_KEY:
        try:
            url = f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}"
            async with session.post(
                url,
                json={
                    "text": script,
                    "model_id": "eleven_turbo_v2_5",
                    "voice_settings": {
                        "stability": 0.4,
                        "similarity_boost": 0.8,
                        "style": 0.0,
                        "use_speaker_boost": True,
                    },
                },
                headers={
                    "Accept": "audio/mpeg",
                    "Content-Type": "application/json",
                    "xi-api-key": ELEVENLABS_API_KEY,
                },
                timeout=60,
            ) as response:
                if response.status == 200:
                    content = await response.read()
                    audio_b64 = base64.b64encode(content).decode("utf-8")
                    print(f"[ElevenLabs] slide_{slide_number} ok ({len(content)} bytes)")
                    return f"data:audio/mpeg;base64,{audio_b64}"
                else:
                    text = await response.text()
                    print(f"[ElevenLabs] Error {response.status}: {safe_str(text[:200])} - falling back")
        except Exception as e:
            print(f"[ElevenLabs] Exception: {safe_str(e)} - falling back")

    # Sarvam TTS fallback
    return await async_generate_audio_sarvam_b64(script, slide_number, session)


async def async_generate_audio_sarvam_b64(script: str, slide_number: int, session: aiohttp.ClientSession) -> str | None:
    """Sarvam AI TTS implicitly called asynchronously. Returns a base64 data URI string or None."""
    if not SARVAM_API_KEY:
        print("[Sarvam TTS] No API key -- skipping audio")
        return None

    try:
        async with session.post(
            "https://api.sarvam.ai/text-to-speech",
            headers={
                "api-subscription-key": SARVAM_API_KEY,
                "Authorization": f"Bearer {SARVAM_API_KEY}",
                "Content-Type": "application/json",
            },
            json={
                "inputs": [script.strip()[:500]],  # clean input and limit length
                "target_language_code": "en-IN",
                "speaker": "priya",             # valid bulbul:v3 speaker
                "pace": 1.0,
                "speech_sample_rate": 22050,
                "enable_preprocessing": True,
                "model": "bulbul:v3",
            },
            timeout=60,
        ) as response:
            if response.status == 200:
                data = await response.json()
                audio_b64 = data.get("audios", [None])[0]
                if audio_b64:
                    print(f"[Sarvam TTS] slide_{slide_number} ok ({len(audio_b64)} b64 chars)")
                    return f"data:audio/wav;base64,{audio_b64}"
            else:
                text = await response.text()
                print(f"[Sarvam TTS] Error {response.status}: {safe_str(text[:300])}")
    except Exception as e:
        print(f"[Sarvam TTS] Exception: {safe_str(e)}")

    return None


async def async_generate_audience_questions(slides_data: list[dict], session: aiohttp.ClientSession) -> list[str]:
    """Generate possible audience questions using Sarvam AI asynchronously."""
    summary_parts = []
    for slide in slides_data:
        summary_parts.append(f"Slide {slide['slide_number']}: {slide.get('title', '')}")
        pts = slide.get("points", [])
        if pts:
            summary_parts.append("  Points: " + ", ".join(pts[:3]))

    presentation_summary = "\n".join(summary_parts)

    prompt = (
        "You are an expert audience member. Based on this presentation, generate 5 insightful "
        "questions that an audience member might ask. Format as a numbered list "
        "(1. Question, 2. Question, etc.). Be specific and thought-provoking.\n\n"
        f"Presentation:\n{presentation_summary}\n\nQuestions:"
    )

    try:
        async with session.post(
            "https://api.sarvam.ai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {SARVAM_API_KEY}",
                "Content-Type": "application/json",
            },
            json={
                "model": "sarvam-m",
                "messages": [{"role": "user", "content": prompt}],
                "max_tokens": 400,
            },
            timeout=30,
        ) as response:
            if response.status == 200:
                data = await response.json()
                content = data["choices"][0]["message"]["content"].strip()
                import re
                content = re.sub(r'<think>.*?</think>', '', content, flags=re.DOTALL).strip()
                questions = []
                for line in content.split("\n"):
                    line = line.strip()
                    if line and (line[0].isdigit() or line.startswith("-")):
                        q = line.lstrip("0123456789.-) ").strip()
                        if q:
                            questions.append(q)
                return questions[:7]
            else:
                text = await response.text()
                print(f"[Sarvam Q&A] Error {response.status}: {safe_str(text[:200])}")
    except Exception as e:
        print(f"[Sarvam Q&A] Exception: {safe_str(e)}")

    # Fallback questions
    topic = slides_data[0]["title"] if slides_data else "this topic"
    return [
        f"What are the key challenges in implementing {topic}?",
        "How does this compare to existing approaches in the field?",
        "What is the future direction of this work?",
        "Can you elaborate on the most surprising finding?",
        "What would be the next steps after this presentation?",
    ]

